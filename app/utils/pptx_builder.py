from pptx import Presentation as PptxPresentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from app.models.presentation_models import PresentationData, Slide, SlideLayout
from app.services.template_service import Template
from app.core.config import logger

def create_presentation_file(data: PresentationData, config, template: Template) -> str:
    """Generates a .pptx file """
    prs = PptxPresentation()

    # 1. Set Aspect Ratio ( Apply a consistent styling.)
    if config.aspect_ratio == "16:9":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
    else: # 4:3
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # 2. Extract colors from template
    bg_color = RGBColor.from_string(template.colors.background)
    text_color = RGBColor.from_string(template.colors.text)
    title_color = RGBColor.from_string(template.colors.title)
    
    for slide_data in data.slides:
        slide_layout_idx = _get_pptx_layout(slide_data.type)
        slide_layout = prs.slide_layouts[slide_layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # 3. Apply background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # 4. Populate slide content with templated styling
        _populate_slide_content(slide, slide_data, template, text_color, title_color)

    # Generate title for filename
    safe_title = "".join(c for c in data.title if c.isalnum() or c in (' ', '_')).rstrip()
    file_path = f"generated_presentations/{safe_title.replace(' ', '_').lower()}.pptx"
    prs.save(file_path)
    logger.info(f"Presentation saved to {file_path}")
    return file_path

def _get_pptx_layout(layout_type: SlideLayout) -> int:
    if layout_type == SlideLayout.TITLE: return 0
    if layout_type == SlideLayout.TWO_COLUMN: return 3
    return 1 # Default to 'Title and Content'

def _populate_slide_content(slide, slide_data: Slide, template: Template, text_color, title_color):
    """
        Method to populates a single slide with content and applies styling from the template.
    """
    # 1. Populate the title (most layouts will have a title)
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = slide_data.title or ""
        # Apply font and color from the template
        font = title_shape.text_frame.paragraphs[0].font
        font.color.rgb = title_color
        font.name = template.font

    # 2. Populate the rest of the content based on slide type and count
    if slide_data.type == SlideLayout.TITLE:
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.subtitle or ""
            font = subtitle.text_frame.paragraphs[0].font
            font.color.rgb = text_color # Use main text color for subtitle
            font.name = template.font

    elif slide_data.type == SlideLayout.BULLET_POINTS:
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()  # Remove any default text

            for point in (slide_data.points or []):
                p = tf.add_paragraph()
                p.text = point
                p.font.color.rgb = text_color
                p.font.name = template.font
                p.level = 0